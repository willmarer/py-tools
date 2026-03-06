import os
import tempfile
from pathlib import Path

from pptx import Presentation

try:
    import pythoncom
    import win32com.client
except ImportError:
    pythoncom = None
    win32com = None


class PPTProcessor:
    def __init__(self, translator, log_callback=None):
        self.translator = translator
        self.log = log_callback or (lambda msg: None)

    def translate_directory(self, input_dir, output_dir):
        input_path = Path(input_dir)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        all_files = []
        all_files.extend(sorted(input_path.glob("*.pptx")))
        all_files.extend(sorted(input_path.glob("*.ppt")))

        stats = {
            "total_files": len(all_files),
            "success_files": 0,
            "failed_files": 0,
            "skipped_files": 0,
            "total_slides": 0,
            "translated_items": 0,
            "errors": [],
        }

        for idx, file_path in enumerate(all_files, start=1):
            self.log("-" * 60)
            self.log(f"[FILE] ({idx}/{len(all_files)}) 开始处理: {file_path.name}")

            try:
                if file_path.suffix.lower() == ".pptx":
                    translated_count, slide_count, save_path = self.translate_pptx(file_path, output_path)
                elif file_path.suffix.lower() == ".ppt":
                    translated_count, slide_count, save_path = self.translate_legacy_ppt(file_path, output_path)
                else:
                    self.log(f"[SKIP] 不支持的文件格式: {file_path.name}")
                    stats["skipped_files"] += 1
                    continue

                stats["success_files"] += 1
                stats["total_slides"] += slide_count
                stats["translated_items"] += translated_count

                self.log(f"[DONE] 保存成功: {save_path}")
                self.log(f"[DONE] 页数: {slide_count}，翻译文本数: {translated_count}")

            except Exception as e:
                stats["failed_files"] += 1
                error_msg = f"{file_path.name} 处理失败: {e}"
                stats["errors"].append(error_msg)
                self.log(f"[ERROR] {error_msg}")

        return stats

    def translate_legacy_ppt(self, ppt_file_path, output_dir):
        """
        处理老格式 .ppt：
        1. 先转成临时 .pptx
        2. 再按 pptx 翻译
        """
        self.log(f"[INFO] 检测到 .ppt 文件，准备自动转换: {ppt_file_path.name}")

        temp_pptx_path = None
        try:
            temp_pptx_path = self.convert_ppt_to_pptx(ppt_file_path)
            self.log(f"[INFO] 转换成功: {Path(temp_pptx_path).name}")

            translated_count, slide_count, final_save_path = self.translate_pptx(
                Path(temp_pptx_path),
                output_dir,
                output_name=f"{ppt_file_path.stem}_translated.pptx"
            )

            return translated_count, slide_count, final_save_path

        finally:
            if temp_pptx_path and os.path.exists(temp_pptx_path):
                try:
                    os.remove(temp_pptx_path)
                    self.log(f"[INFO] 已清理临时文件: {Path(temp_pptx_path).name}")
                except Exception as cleanup_error:
                    self.log(f"[WARN] 清理临时文件失败: {cleanup_error}")

    def convert_ppt_to_pptx(self, ppt_file_path):
        """
        使用本机 PowerPoint 把 .ppt 转成临时 .pptx
        前提：Windows + 已安装 Microsoft PowerPoint
        """
        if win32com is None or pythoncom is None:
            raise RuntimeError("缺少 pywin32，请先执行: pip install pywin32")

        pythoncom.CoInitialize()
        powerpoint = None
        presentation = None

        try:
            try:
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            except Exception:
                raise RuntimeError("无法调用 PowerPoint。请确认本机已安装 Microsoft PowerPoint。")

            powerpoint.Visible = 1

            ppt_file_path = str(Path(ppt_file_path).resolve())

            temp_dir = tempfile.gettempdir()
            temp_pptx_path = os.path.join(
                temp_dir,
                f"ppt_translate_temp_{Path(ppt_file_path).stem}.pptx"
            )

            # 如果旧临时文件存在，先删掉
            if os.path.exists(temp_pptx_path):
                try:
                    os.remove(temp_pptx_path)
                except Exception:
                    pass

            self.log("[INFO] 正在调用 PowerPoint 转换 .ppt -> .pptx ...")

            # WithWindow=False
            presentation = powerpoint.Presentations.Open(ppt_file_path, WithWindow=False)

            # 24 = ppSaveAsOpenXMLPresentation (.pptx)
            presentation.SaveAs(temp_pptx_path, 24)
            presentation.Close()
            presentation = None

            if not os.path.exists(temp_pptx_path):
                raise RuntimeError("PowerPoint 转换失败，未生成 pptx 文件。")

            return temp_pptx_path

        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass

            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except Exception:
                    pass

            pythoncom.CoUninitialize()

    def translate_pptx(self, file_path, output_dir, output_name=None):
        prs = Presentation(str(file_path))
        slide_count = len(prs.slides)
        translated_count = 0

        for slide_index, slide in enumerate(prs.slides, start=1):
            self.log(f"  [SLIDE] {slide_index}/{slide_count}")
            translated_count += self._translate_slide(slide)

            if slide.has_notes_slide:
                try:
                    notes_frame = slide.notes_slide.notes_text_frame
                    translated_count += self._translate_text_frame(notes_frame)
                except Exception:
                    pass

        if not output_name:
            output_name = f"{Path(file_path).stem}_translated.pptx"

        save_path = Path(output_dir) / output_name
        prs.save(str(save_path))

        return translated_count, slide_count, str(save_path)

    def _translate_slide(self, slide):
        translated_count = 0
        for shape in slide.shapes:
            translated_count += self._translate_shape(shape)
        return translated_count

    def _translate_shape(self, shape):
        translated_count = 0

        # 文本框
        if hasattr(shape, "text_frame"):
            translated_count += self._translate_text_frame(shape.text_frame)

        # 表格
        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    translated_count += self._translate_text_frame(cell.text_frame)

        # 组合图形
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                translated_count += self._translate_shape(sub_shape)

        return translated_count

    def _translate_text_frame(self, text_frame):
        translated_count = 0

        if not text_frame:
            return translated_count

        for paragraph in text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)

            if not full_text or not full_text.strip():
                continue

            # 换行转空格，便于短语匹配
            normalized_text = full_text.replace("\r", " ").replace("\n", " ")
            translated = self.translator.translate(normalized_text)

            if translated != normalized_text:
                paragraph.clear()
                paragraph.add_run().text = translated
                translated_count += 1

        return translated_count
